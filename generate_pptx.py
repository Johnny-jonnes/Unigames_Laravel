import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = subtitle_text

def add_image_slide(prs, title_text, content_text, image_path):
    # Layout 5 is title only, we can add image manually
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    # Add text box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = content_text
    tf.word_wrap = True
    
    # Add image if exists
    if os.path.exists(image_path):
        img_left = Inches(1)
        img_top = Inches(2.5)
        # Calculate width to fit the slide (max width ~8 inches, max height ~4.5 inches)
        slide.shapes.add_picture(image_path, img_left, img_top, width=Inches(8))
    else:
        print(f"Warning: Image {image_path} not found.")

def generate_manuel_pptx():
    prs = Presentation()
    
    # Set slide width and height for 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    create_title_slide(prs, "Manuel Utilisateur - UniGames", "Guide Complet d'Utilisation de la Plateforme")
    
    slides_data = [
        ("1. Page de Connexion", "Authentification securisee pour l'acces a la plateforme.", "docs/images/2_login.jpeg"),
        ("2. Tableau de Bord Principal", "Vue d'ensemble avec KPIs, derniers resultats et prochains matchs.", "docs/images/4_dashboard_principal.jpeg"),
        ("3. Liste des Facultes", "Consultation de toutes les facultes inscrites a l'edition.", "docs/images/5_liste_facultes.jpeg"),
        ("4. Details d'une Faculte", "Vue detaillee d'une faculte et de ses equipes associees.", "docs/images/7_details_faculte.jpeg"),
        ("5. Repertoire des Equipes", "Liste exhaustive de toutes les equipes du tournoi.", "docs/images/11_repertoire_equipes.jpeg"),
        ("6. Inscription d'une Equipe", "Formulaire d'ajout d'une nouvelle equipe.", "docs/images/9_ajouter_equipe.jpeg"),
        ("7. Repertoire des Joueurs", "Liste de tous les joueurs inscrits.", "docs/images/12_repertoire_joueurs.jpeg"),
        ("8. Enregistrer un Joueur", "Formulaire pour ajouter un nouveau joueur a une equipe.", "docs/images/10_ajouter_joueur.jpeg"),
        ("9. Calendrier des Rencontres", "Tous les matchs organises par discipline.", "docs/images/13_calendrier_rencontres.jpeg"),
        ("10. Details d'un Match & Score", "Fiche complete du match et formulaire de saisie du score et des buteurs.", "docs/images/14_details_match.jpeg"),
        ("11. Classements et Statistiques", "Classement par discipline et top des meilleurs buteurs.", "docs/images/15_classements_statistiques.jpeg"),
        ("12. Arbre du Tournoi", "Visualisation interactive de l'avancement des phases finales.", "docs/images/18_arbre_tournoi.jpeg"),
        ("13. Gestion du Staff", "Interface d'administration pour la gestion des comptes utilisateurs.", "docs/images/16_gestion_staff.jpeg"),
        ("14. Profil et Parametres", "Gestion du compte personnel et changement de mot de passe.", "docs/images/17_profil_parametres.jpeg")
    ]
    
    for title, text, img in slides_data:
        add_image_slide(prs, title, text, img)
        
    prs.save("docs/8_Manuel_Utilisateur_Presentation.pptx")
    print("[OK] 8_Manuel_Utilisateur_Presentation.pptx")

def generate_maquette_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    create_title_slide(prs, "Dossier de Maquettage - UniGames", "Interfaces Utilisateur et Design System")
    
    slides_data = [
        ("1. Interface de Connexion", "Design epure avec mise en valeur du formulaire et gestion des erreurs en temps reel.", "docs/images/2_login.jpeg"),
        ("2. Dashboard & KPIs", "Tableau de bord utilisant des cartes (Enterprise Cards) pour une lisibilite optimale des donnees cles.", "docs/images/4_dashboard_principal.jpeg"),
        ("3. Calendrier et Accordeons", "Utilisation de composants repliables pour organiser les matchs par discipline sportive.", "docs/images/13_calendrier_rencontres.jpeg"),
        ("4. Fiche Match & Scoreboard", "Mise en page centralisee avec un scoreboard visuel XXL et badges de statut.", "docs/images/14_details_match.jpeg"),
        ("5. Visualisation de l'Arbre", "Design horizontal scrollable avec cartes stylisees selon les phases (ex: carte doree pour la finale).", "docs/images/18_arbre_tournoi.jpeg"),
        ("6. Tableaux de Classement", "Design de tableau data-grid avec alternance de couleurs (zebra) pour faciliter la lecture.", "docs/images/15_classements_statistiques.jpeg")
    ]
    
    for title, text, img in slides_data:
        add_image_slide(prs, title, text, img)
        
    prs.save("docs/9_Maquette_Presentation.pptx")
    print("[OK] 9_Maquette_Presentation.pptx")

if __name__ == "__main__":
    print("Generation des PowerPoint en cours...")
    if not os.path.exists("docs"):
        os.makedirs("docs")
    generate_manuel_pptx()
    generate_maquette_pptx()
    print("Termine avec succes !")
